#!/usr/bin/env python3
"""
Créer un PPTX simple avec structure minimale garantie.
"""
import os
import sys

# Essayer d'importer python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    USE_PPTX = True
except ImportError:
    USE_PPTX = False
    print("python-pptx non disponible, utilisation de méthode ZIP brute")

def create_with_pptx():
    """Créer avec python-pptx."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slides_data = [
        ("Cat Clicker", "Un jeu de clic addictif"),
        ("Concept du Jeu", "Cliquez sur le chat pour augmenter votre score"),
        ("Caractéristiques", "7 niveaux progressifs\n2 thèmes de couleurs\nEffets sonores\nSystème de vies\nScore persistant"),
        ("Mécaniques", "Clic de souris\nProgression automatique\nMultiplicateurs de points\nLimite de temps"),
        ("Technologie", "Développé en C avec SDL2\nGraphismes haute performance\nGestion audio intégrée"),
        ("Objectifs", "Atteindre le niveau 7\nMaximiser votre score\nCompléter tous les thèmes"),
        ("Merci!", "Amusez-vous bien avec Cat Clicker"),
    ]
    
    for title, content in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        content_shape.text = content
    
    output_path = r"c:\Users\amity\OneDrive\Desktop\cat game\Cat_Clicker_Presentation.pptx"
    prs.save(output_path)
    return output_path

def create_manual_pptx():
    """Créer avec ZIP et XML brut."""
    import zipfile
    
    output_path = r"c:\Users\amity\OneDrive\Desktop\cat game\Cat_Clicker_Presentation.pptx"
    
    with zipfile.ZipFile(output_path, 'w') as pptx:
        # [Content_Types].xml - MINIMAL
        pptx.writestr('[Content_Types].xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
<Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slides/slide2.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slides/slide3.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slides/slide4.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slides/slide5.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slides/slide6.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slides/slide7.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
</Types>''')
        
        pptx.writestr('_rels/.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
</Relationships>''')
        
        pptx.writestr('docProps/core.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
<dc:creator>User</dc:creator>
<dc:title>Cat Clicker</dc:title>
</cp:coreProperties>''')
        
        pptx.writestr('ppt/_rels/presentation.xml.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
<Relationship Id="rId4" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
<Relationship Id="rId5" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide2.xml"/>
<Relationship Id="rId6" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide3.xml"/>
<Relationship Id="rId7" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide4.xml"/>
<Relationship Id="rId8" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide5.xml"/>
<Relationship Id="rId9" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide6.xml"/>
<Relationship Id="rId10" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide7.xml"/>
<Relationship Id="rId11" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="theme/theme1.xml"/>
</Relationships>''')
        
        pptx.writestr('ppt/presentation.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<p:sldIdLst>
<p:sldId id="256" r:id="rId4"/>
<p:sldId id="257" r:id="rId5"/>
<p:sldId id="258" r:id="rId6"/>
<p:sldId id="259" r:id="rId7"/>
<p:sldId id="260" r:id="rId8"/>
<p:sldId id="261" r:id="rId9"/>
<p:sldId id="262" r:id="rId10"/>
</p:sldIdLst>
</p:presentation>''')
        
        # Créer les slides
        slides_data = [
            ("Cat Clicker", "Un jeu de clic addictif"),
            ("Concept du Jeu", "Cliquez sur le chat pour augmenter votre score"),
            ("Caracteristiques", "7 niveaux progressifs, 2 themes, Effets sonores, Systeme de vies, Score persistant"),
            ("Mecaniques", "Clic de souris, Progression automatique, Multiplicateurs de points, Limite de temps"),
            ("Technologie", "Developpe en C avec SDL2, Graphismes haute performance, Gestion audio integree"),
            ("Objectifs", "Atteindre le niveau 7, Maximiser votre score, Completer tous les themes"),
            ("Merci", "Amusez-vous bien avec Cat Clicker"),
        ]
        
        for i, (title, content) in enumerate(slides_data, 1):
            # Remplacer les caractères spéciaux
            title_safe = title.replace('é', 'e').replace('à', 'a').replace('â', 'a')
            content_safe = content.replace('é', 'e').replace('à', 'a').replace('â', 'a')
            
            slide_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
<p:cSld>
<p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name="Title 1"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="9144000" cy="6858000"/><a:chOff x="0" y="0"/><a:chExt cx="9144000" cy="6858000"/></a:xfrm></p:grpSpPr>
<p:sp>
<p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
<p:spPr><a:xfrm><a:off x="457200" y="274638"/><a:ext cx="8229600" cy="1143000"/></a:xfrm></p:spPr>
<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{title_safe}</a:t></a:r></a:p></p:txBody>
</p:sp>
<p:sp>
<p:nvSpPr><p:cNvPr id="3" name="Content"/><p:cNvSpPr><p:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr/></p:nvSpPr>
<p:spPr><a:xfrm><a:off x="457200" y="1417638"/><a:ext cx="8229600" cy="5016000"/></a:xfrm></p:spPr>
<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{content_safe}</a:t></a:r></a:p></p:txBody>
</p:sp>
</p:spTree>
</p:cSld>
</p:sld>'''
            pptx.writestr(f'ppt/slides/slide{i}.xml', slide_xml)
        
        # slideLayout
        pptx.writestr('ppt/slideLayouts/slideLayout1.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name="Title Slide Layout"/><p:cNvGrpSpPr/><p:nvPr><p:ph/></p:nvPr></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="9144000" cy="6858000"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
</p:sldLayout>''')
        
        # slideMaster
        pptx.writestr('ppt/slideMasters/slideMaster1.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
<p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name="Slide Master"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="9144000" cy="6858000"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
<p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
</p:sldMaster>''')
        
        # slideMaster rels
        pptx.writestr('ppt/slideMasters/_rels/slideMaster1.xml.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>''')
        
        # slideLayout rels
        pptx.writestr('ppt/slideLayouts/_rels/slideLayout1.xml.rels', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>''')
        
        # theme
        pptx.writestr('ppt/theme/theme1.xml', '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Office Theme">
<a:themeElements><a:clrScheme name="Office"><a:dk1><a:srgbClr val="000000"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1></a:clrScheme></a:themeElements>
</a:theme>''')
    
    return output_path

if __name__ == '__main__':
    try:
        if USE_PPTX:
            path = create_with_pptx()
            print(f"✅ PPTX créé avec python-pptx: {path}")
        else:
            path = create_manual_pptx()
            print(f"✅ PPTX créé manuellement: {path}")
    except Exception as e:
        print(f"❌ Erreur: {e}")
        sys.exit(1)
